"""
File processors for extracting text from various file formats.
Supports PDF, DOCX, PPTX, and TXT files.
"""

import os
from typing import Dict, Optional
from io import BytesIO
import PyPDF2
from docx import Document
from pptx import Presentation


class FileProcessor:
    """Process various file formats and extract text."""
    
    @staticmethod
    def process_file(file_path: str = None, file_content: bytes = None, 
                    file_name: str = None, mime_type: str = None) -> Dict:
        """
        Process a file and extract text content.
        
        Args:
            file_path: Path to file (if available)
            file_content: File content as bytes
            file_name: Name of the file
            mime_type: MIME type of the file
            
        Returns:
            Dictionary with 'text', 'metadata', and 'error' fields
        """
        if not file_path and not file_content:
            return {'text': '', 'metadata': {}, 'error': 'No file provided'}
        
        # Determine file type
        if file_name:
            ext = os.path.splitext(file_name)[1].lower()
        elif file_path:
            ext = os.path.splitext(file_path)[1].lower()
        else:
            ext = FileProcessor._get_extension_from_mime(mime_type)
        
        # Read file content if path provided
        if file_path and not file_content:
            try:
                with open(file_path, 'rb') as f:
                    file_content = f.read()
            except Exception as e:
                return {'text': '', 'metadata': {'error': str(e)}, 'error': str(e)}
        
        # Process based on extension
        try:
            if ext == '.pdf':
                return FileProcessor._process_pdf(file_content, file_name)
            elif ext in ['.docx', '.doc']:
                return FileProcessor._process_docx(file_content, file_name)
            elif ext in ['.pptx', '.ppt']:
                return FileProcessor._process_pptx(file_content, file_name)
            elif ext == '.txt':
                return FileProcessor._process_txt(file_content, file_name)
            else:
                return {
                    'text': '',
                    'metadata': {'file_name': file_name, 'extension': ext},
                    'error': f'Unsupported file type: {ext}'
                }
        except Exception as e:
            return {
                'text': '',
                'metadata': {'file_name': file_name},
                'error': f'Error processing file: {str(e)}'
            }
    
    @staticmethod
    def _process_pdf(content: bytes, file_name: str = None) -> Dict:
        """Extract text from PDF file."""
        try:
            pdf_file = BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_parts = []
            num_pages = len(pdf_reader.pages)
            
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    text = page.extract_text()
                    if text:
                        text_parts.append(f"--- Page {page_num + 1} ---\n{text}")
                except Exception as e:
                    print(f"Error extracting page {page_num + 1}: {str(e)}")
            
            full_text = "\n\n".join(text_parts)
            
            return {
                'text': full_text,
                'metadata': {
                    'file_name': file_name,
                    'file_type': 'pdf',
                    'num_pages': num_pages
                },
                'error': None
            }
        except Exception as e:
            return {
                'text': '',
                'metadata': {'file_name': file_name, 'file_type': 'pdf'},
                'error': f'PDF processing error: {str(e)}'
            }
    
    @staticmethod
    def _process_docx(content: bytes, file_name: str = None) -> Dict:
        """Extract text from DOCX file."""
        try:
            doc_file = BytesIO(content)
            doc = Document(doc_file)
            
            text_parts = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                if table_text:
                    text_parts.append("\n".join(table_text))
            
            full_text = "\n\n".join(text_parts)
            
            return {
                'text': full_text,
                'metadata': {
                    'file_name': file_name,
                    'file_type': 'docx',
                    'num_paragraphs': len(doc.paragraphs)
                },
                'error': None
            }
        except Exception as e:
            return {
                'text': '',
                'metadata': {'file_name': file_name, 'file_type': 'docx'},
                'error': f'DOCX processing error: {str(e)}'
            }
    
    @staticmethod
    def _process_pptx(content: bytes, file_name: str = None) -> Dict:
        """Extract text from PPTX file."""
        try:
            pptx_file = BytesIO(content)
            prs = Presentation(pptx_file)
            
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                slide_text.append(f"--- Slide {slide_num + 1} ---")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # More than just the header
                    text_parts.append("\n".join(slide_text))
            
            full_text = "\n\n".join(text_parts)
            
            return {
                'text': full_text,
                'metadata': {
                    'file_name': file_name,
                    'file_type': 'pptx',
                    'num_slides': len(prs.slides)
                },
                'error': None
            }
        except Exception as e:
            return {
                'text': '',
                'metadata': {'file_name': file_name, 'file_type': 'pptx'},
                'error': f'PPTX processing error: {str(e)}'
            }
    
    @staticmethod
    def _process_txt(content: bytes, file_name: str = None) -> Dict:
        """Extract text from TXT file."""
        try:
            # Try UTF-8 first, fallback to latin-1
            try:
                text = content.decode('utf-8')
            except UnicodeDecodeError:
                text = content.decode('latin-1')
            
            return {
                'text': text,
                'metadata': {
                    'file_name': file_name,
                    'file_type': 'txt'
                },
                'error': None
            }
        except Exception as e:
            return {
                'text': '',
                'metadata': {'file_name': file_name, 'file_type': 'txt'},
                'error': f'TXT processing error: {str(e)}'
            }
    
    @staticmethod
    def _get_extension_from_mime(mime_type: str) -> str:
        """Get file extension from MIME type."""
        mime_map = {
            'application/pdf': '.pdf',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
            'text/plain': '.txt'
        }
        return mime_map.get(mime_type, '')

